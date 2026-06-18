# src/app.py
import datetime
import os
import re
import subprocess
import time
import sys
import json
import threading
import uuid
import openai  
import base64
import requests
import shutil
import zlib
import traceback 
import urllib.parse     # import missing 해결

from util.date_query import run_date_range_query

from dotenv import load_dotenv
from flask import Flask, request, jsonify, send_from_directory, Response, stream_with_context
from flask_cors import CORS
import fitz  # PyMuPDF
from docx import Document
import olefile
import csv
from pptx import Presentation
from openpyxl import load_workbook
from flask import send_from_directory

# Job 이용 공통함수 import
from util.jobs.job_store import *
from util.jobs.job_run import start_graph_pipeline_background, start_graph_update_pipeline_background
from config.settings import *
from util.user_path import UserPaths
from util.database.db_reader import get_mail_stats, get_keyword_stats, get_mail_sync_stats, get_user_rating_stats, get_high_affinity_person_stats, get_mail_date_range, get_mail_exchange_stats, get_date_range_person_stats, get_keywords_by_person_date

from util.database.db_writer import (
    save_query_to_db,
    init_processed_attachments_table,
    init_keyword_mail_table,
    filter_unprocessed_attachments,
    mark_attachments_as_processed
)
from util.extract_statics import start_statics_pipeline_background

from util.sse_broadcaster import subscribe, unsubscribe

from config.db import get_db_connection

# 환경변수 로드
load_dotenv("src/parquet/.env")

# Flask 앱 초기화
app = Flask(__name__)
CORS(app)


# 서버 시작 시 테이블 초기화 실행
init_processed_attachments_table()
init_keyword_mail_table()

# Apps Script Web App URL
WEBAPP_URL = "https://script.google.com/macros/s/AKfycbwAk_JabdKuGUHIVcaKeEnY1DUiYb0uqkiu-KdUG67Zf1U3D8k-F06RGS5043k_fZS8MQ/exec"


# 한글 출력 시 깨지거나 에러 나는 것 방지
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
if hasattr(sys.stderr, "reconfigure"):
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")



# 유틸 함수

# GraphRAG CLI 실행
def _run_graphrag(message, resMethod, raw_message, paths, resType):

    def decode_output(b: bytes) -> str:
        if not b:
            return ""
        for enc in ("utf-8", "cp949", "euc-kr"):
            try:
                return b.decode(enc)
            except UnicodeDecodeError:
                pass
        return b.decode("utf-8", errors="replace")

    python_command = [
        'graphrag', 'query',
        '--root', paths.GRAPHRAG_ROOT,
        '--response-type', resType,
        '--method', resMethod,
        '--query', message
    ]

    start_time = time.time()

    result = subprocess.run(
        python_command,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        env=os.environ.copy(),
        text=False
    )
    elapsed = time.time() - start_time
    print(f'execution_time : {elapsed}')
    try:
        save_query_to_db(paths.GMAIL_ID, raw_message, elapsed, resMethod)
    except Exception as e:
        print(f"[WARN] query DB 저장 실패 (무시): {e}")

    stdout_text = decode_output(result.stdout)
    stderr_text = decode_output(result.stderr)

    if result.returncode != 0:
        raise RuntimeError(stderr_text or stdout_text or 'GraphRAG 실행 오류')

    print(stdout_text)

    match = re.search(r'SUCCESS: (?:Local|Global) Search Response:\s*(.*)', stdout_text, re.DOTALL)
    answer = match.group(1).strip() if match else stdout_text.strip()

    answer = re.sub(r'\[Data:.*?\]|\[데이터:.*?\]', '', answer)
    answer = re.sub(r'\*+|#+', '', answer)
    answer = answer.strip()
    print(answer)
    return answer.strip()

# 텍스트 → 캘린더 JSON 변환
def _convert_to_calendar_json(text):
    client = openai.OpenAI(api_key=os.environ.get("GRAPHRAG_API_KEY"))
    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            response_format={"type": "json_object"},
            messages=[
                {
                    "role": "system",
                    "content": (
                        "너는 이메일 내용을 분석해서 캘린더 일정을 추출하는 도우미야."
                        "날짜/시간/일정 정보를 추출해서 반드시 JSON으로만 응답해. "
                        "이메일의 제목과 본문을 함께 분석해서 캘린더에 적합한 새로운 일정 제목(title)을 만들어."
                        "메일 제목을 그대로 복사하지 말고, 실제 일정의 목적이 드러나도록 자연스럽고 짧게 작성해."
                        "예를 들면 '회의 안내' 같은 제목이 있더라도, 본문이 캡스톤 발표 회의에 대한 내용이면 title는 '캡스톤 발표 회의'처럼 만들어."
                        "title은 5~20자 정도의 짧고 명확한 한국어로 작성해."
                        "description은 일정과 관련된 핵심 내용을 간단히 넣어"
                        "형식: {\"events\": [{\"title\": \"제목\", \"startTime\": \"2026-02-26 Time 09:00:00\", "
                        "\"endTime\": \"2026-02-26 Time 10:00:00\", \"description\": \"\"}]} "
                        "일정 없으면 {\"events\": []}"
                    )
                },
                {"role": "user", "content": text}
            ]
        )
        return json.loads(response.choices[0].message.content)
    except Exception as e:
        print(f"[calendar convert error] {e}")
        return {"events": []}

# 첨부파일 텍스트 요약
def _summarize_attachment(text: str, filename: str) -> str:
    pure_len = len(text.replace(" ", "").replace("\n", ""))
    if pure_len < 500:
        return text

    prompt_path = os.path.join("parquet_template", "prompts", "summarize_attachment.txt")
    with open(prompt_path, "r", encoding="utf-8") as f:
        prompt = f.read().strip()

    client = openai.OpenAI(api_key=os.environ.get("GRAPHRAG_API_KEY"))
    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": prompt},
                {"role": "user", "content": f"파일명: {filename}\n\n{text}"}
            ],
            max_tokens=150
        )
        result = response.choices[0].message.content.strip()
        REFUSAL_PREFIXES = ("죄송", "I'm sorry", "I'm unable", "I cannot", "Sorry")
        if result.startswith(REFUSAL_PREFIXES):
            print(f"[summarize_attachment] LLM 거부 응답 감지: {filename}")
            return ""
        return result
    except Exception as e:
        print(f"[summarize_attachment error] {e}")
        return ""

# PDF 파일에서 텍스트 추출
def _extract_text_from_pdf(file_path):
    text = ""
    try:
        doc = fitz.open(file_path)
        for page in doc:
            text += page.get_text()
        doc.close()
    except Exception as e:
        print(f"[PDF Extract Error] {e}")
    return text

# Word 파일에서 텍스트 추출
def _extract_text_from_docx(file_path):
    text = ""
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        print(f"[Docx Extract Error] {e}")
    return text

# HWP 파일에서 텍스트 추출
def _extract_text_from_hwp(file_path):
    text = ""
    try:
        f = olefile.OleFileIO(file_path)
        dirs = f.listdir()
        sections = [d for d in dirs if "BodyText/Section" in "/".join(d)]
        for section in sections:
            stream = f.openstream("/".join(section))
            data = stream.read()
            try:
                decompressed = zlib.decompress(data, -15)
                decoded_text = decompressed.decode("utf-16-le", errors="ignore")
                clean_text = "".join(c for c in decoded_text if c.isalnum() or c in " \n\t.,()[]")
                text += clean_text + "\n"
            except Exception as e:
                print(f"[HWP Decode Error in {section}] {e}")
        f.close()
    except Exception as e:
        print(f"[HWP Extract Error] {e}")
    return text

# TXT 파일에서 텍스트 추출
def _extract_text_from_txt(file_path):
    text = ""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
    except UnicodeDecodeError:
        try:
            with open(file_path, "r", encoding="cp949") as f:
                text = f.read()
        except Exception as e:
            print(f"[TXT Extract Error] {e}")
    except Exception as e:
        print(f"[TXT Extract Error] {e}")
    return text

# PPTX 파일에서 텍스트 추출
def _extract_text_from_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
    except Exception as e:
        print(f"[PPTX Extract Error] {e}")
    return text

# XLSX 파일에서 텍스트 추출
def _extract_text_from_xlsx(file_path):
    text = ""
    try:
        wb = load_workbook(file_path, data_only=True)
        for ws in wb.worksheets:
            text += f"[Sheet] {ws.title}\n"
            for row in ws.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                if any(v.strip() for v in row_values):
                    text += " | ".join(row_values) + "\n"
            text += "\n"
    except Exception as e:
        print(f"[XLSX Extract Error] {e}")
    return text

# CSV 파일에서 텍스트 추출
def _extract_text_from_csv(file_path):
    text = ""
    try:
        with open(file_path, "r", encoding="utf-8", newline="") as f:
            reader = csv.reader(f)
            for row in reader:
                row_values = [str(cell) if cell is not None else "" for cell in row]
                text += " | ".join(row_values) + "\n"
    except UnicodeDecodeError:
        try:
            with open(file_path, "r", encoding="cp949", newline="") as f:
                reader = csv.reader(f)
                for row in reader:
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    text += " | ".join(row_values) + "\n"
        except Exception as e:
            print(f"[CSV Extract Error] {e}")
    except Exception as e:
        print(f"[CSV Extract Error] {e}")
    return text

# 파일명에서 경로/위험 문자 제거
def _sanitize_filename(name: str) -> str:
    name = os.path.basename(name or "attachment.bin").strip()
    name = re.sub(r"[^A-Za-z0-9._-]", "_", name)
    return name or "attachment.bin"

# attachment payload에서 base64를 받아 서버 로컬에 파일 저장
def _save_attachment_from_base64(file_info: dict, save_dir: str) -> tuple[str, str]:
    original_name = file_info.get("name") or "attachment.bin"
    safe_name = _sanitize_filename(original_name)
    mail_id = str(file_info.get("mail_id") or "no_mail_id")
    data_base64 = file_info.get("data_base64") or ""

    if not data_base64:
        raise ValueError(f"attachment data_base64 missing: {original_name}")

    os.makedirs(save_dir, exist_ok=True)

    ext = os.path.splitext(safe_name)[1].lower()
    unique_name = f"{mail_id}_{uuid.uuid4().hex[:8]}{ext or '.bin'}"
    saved_path = os.path.join(save_dir, unique_name)

    if "," in data_base64 and "base64" in data_base64[:100]:
        data_base64 = data_base64.split(",", 1)[1]

    file_bytes = base64.b64decode(data_base64)

    with open(saved_path, "wb") as f:
        f.write(file_bytes)

    return saved_path, original_name

# 메일 블록에서 'ID: ...' 값을 추출
def _extract_mail_id_from_block(block: str) -> str | None:
    m = re.search(r"^\s*ID:\s*(.+?)\s*$", block, re.MULTILINE)
    return m.group(1).strip() if m else None

# mail_id 기준으로 첨부 텍스트를 각 메일 블록 하단에 삽입한 후 다시 append
def _merge_attachments_into_mail_blocks(content: str, attachment_texts_by_mail: dict[str, list[dict]]) -> str:
    parts = content.split(MAIL_BLOCK_SEP)
    merged_blocks = []

    for part in parts:
        block = part.strip()
        if not block:
            continue

        block_text = f"{MAIL_BLOCK_SEP}\n{block}\n{MAIL_BLOCK_SEP}"

        mail_id = _extract_mail_id_from_block(block_text)
        if not mail_id:
            merged_blocks.append(block_text)
            continue

        attachment_entries = attachment_texts_by_mail.get(mail_id, [])
        if not attachment_entries:
            merged_blocks.append(block_text)
            continue

        attachment_section = "\n[첨부 추출 내용]\n"
        for item in attachment_entries:
            attachment_section += f"[File name] {item['name']}\n{item['text']}\n"

        insert_pos = block_text.rfind(MAIL_BLOCK_SEP)
        if insert_pos == -1:
            merged_blocks.append(block_text + attachment_section)
        else:
            merged_blocks.append(
                block_text[:insert_pos].rstrip() + "\n\n" +
                attachment_section.rstrip() + "\n" +
                MAIL_BLOCK_SEP
            )

    return "\n".join(merged_blocks) + "\n"

# 텍스트에서 메일별로 구분
def _split_mail_blocks(text):
    parts = text.split(MAIL_BLOCK_SEP)
    blocks = []

    for p in parts:
        p = p.strip()
        if not p:
            continue
        block = MAIL_BLOCK_SEP + "\n" + p
        if not block.endswith(MAIL_BLOCK_SEP):
            block += "\n" + MAIL_BLOCK_SEP
        blocks.append(block)

    return blocks

def _renumber_mail_blocks(text: str) -> str:
    blocks = _split_mail_blocks(text)
    result = []
    for i, block in enumerate(blocks, start=1):
        renumbered = re.sub(r'\[메일 \d+\]', f'[메일 {i}]', block)
        result.append(renumbered)
    return "\n".join(result) + "\n"

# 메일 id들 추출해서 집합으로 반환
def _extract_message_ids(text):
    return set(re.findall(r"^\s*ID:\s*(.+?)\s*$", text, flags=re.MULTILINE))

# 메일 블록에서 "날짜:" 부분 파싱해서 datetime 객체로 반환
def _extract_block_for_sort(block):
    for line in block.splitlines():
        if line.startswith("날짜:"):
            raw = line.replace("날짜:", "").strip()
            try:
                return datetime.datetime.strptime(raw, "%Y-%m-%d %H:%M:%S")
            except Exception:
                return datetime.datetime.min
    return datetime.datetime.min

# 현재 mail_latest.txt 파일 전체 문자열로 읽어서 반환
def _read_latest_text(paths):
    if not os.path.exists(paths.MAIL_LATEST_PATH):
        return ""
    with open(paths.MAIL_LATEST_PATH, "r", encoding="utf-8") as f:
        return f.read()

# 업데이트 시 생기는 input 폴더 속 새로운 메일 증분 파일 삭제
def _delete_incremental_files(paths):
    os.makedirs(paths.MAIL_DIR, exist_ok=True)

    for name in os.listdir(paths.MAIL_DIR):
        is_inc_txt = name.startswith("inc_") and name.endswith(".txt")
        is_inc_csv = name.startswith("inc_") and name.endswith(".csv")
        is_att_txt = name == "attachment_latest.txt"

        if is_inc_txt or is_inc_csv or is_att_txt:
            path = os.path.join(paths.MAIL_DIR, name)
            try:
                os.remove(path)
            except Exception as e:
                print(f"[UPLOAD] failed to remove incremental file: {path} / {e}")

# 업데이트 시 생기는 update_output 폴더 속 새로운 결과 파일 삭제
def _delete_old_update_files(paths):
    update_output_dir = paths.UPDATE_DIR
    if not os.path.exists(update_output_dir):
        return

    folders = sorted([
        f for f in os.listdir(update_output_dir)
        if os.path.isdir(os.path.join(update_output_dir, f))
    ])

    for folder in folders[:-1]:
        folder_path = os.path.join(update_output_dir, folder)
        try:
            shutil.rmtree(folder_path)
            print(f"[CLEANUP] 삭제: {folder_path}")
        except Exception as e:
            print(f"[CLEANUP] 삭제 실패 (무시): {e}")

# 증분 파일 저장경로 생성
def _build_incremental_path(filename: str, paths) -> str:
    safe_name = _sanitize_filename(filename or "")
    if not safe_name.startswith("inc_"):
        safe_name = f"inc_{datetime.datetime.now().strftime('%Y-%m-%d_%H%M%S')}.txt"
    return os.path.join(paths.MAIL_DIR, safe_name)

# json 파일 읽어서 dict로 파싱 후 반환
def _read_json_file(path):
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)

# 인덱스 여부 확인
def _is_index_ready(paths):
    stats_path = os.path.join(paths.GRAPHRAG_ROOT, "output", "stats.json")

    try:
        required_paths = [paths.MAIL_LATEST_PATH, stats_path]

        for path in required_paths:
            if not os.path.exists(path):
                print(f"[INDEX READY] missing: {path}")
                return False
            if os.path.getsize(path) == 0:
                print(f"[INDEX READY] empty file: {path}")
                return False

        _read_json_file(stats_path)
        return True

    except Exception as e:
        print(f"[INDEX READY] invalid index state: {e}")
        return False

# 백그라운드: 첨부파일 텍스트 추출 → 요약 → attachment_latest.txt 저장 → graphrag update
def _run_attachment_pipeline(job_id: str, paths, attachments: list, env: dict, is_last):
    from util.jobs.job_run import build_graphrag_update, build_graph_json

    print(f"[JOB][attachment] START job_id={job_id}")
    update_job(job_id, status="running", progress=0, message="첨부파일 텍스트 추출 중")

    try:
        attachment_texts_by_mail: dict[str, list[dict]] = {}

        # 1) 첨부파일 저장 + 텍스트 추출
        for file_info in attachments:
            f_name = file_info.get("name") or "attachment.bin"
            mime = (file_info.get("mime") or "").lower()
            mail_id = str(file_info.get("mail_id") or "").strip()

            if not mail_id:
                continue

            try:
                saved_path, original_name = _save_attachment_from_base64(file_info, paths.ATTACHMENT_DIR)
                ext = os.path.splitext(original_name)[-1].lower()
                file_text = ""

                if ext == ".pdf" or "pdf" in mime:     file_text = _extract_text_from_pdf(saved_path)
                elif ext == ".docx":                    file_text = _extract_text_from_docx(saved_path)
                elif ext == ".hwp":                     file_text = _extract_text_from_hwp(saved_path)
                elif ext == ".txt" or "plain" in mime:  file_text = _extract_text_from_txt(saved_path)
                elif ext == ".pptx":                    file_text = _extract_text_from_pptx(saved_path)
                elif ext == ".xlsx":                    file_text = _extract_text_from_xlsx(saved_path)
                elif ext == ".csv":                     file_text = _extract_text_from_csv(saved_path)

                if file_text and file_text.strip():
                    attachment_texts_by_mail.setdefault(mail_id, []).append({
                        "name": original_name,
                        "text": file_text.strip()
                    })

            except Exception as e:
                print(f"[JOB][attachment] extract error {f_name}: {e}")

        update_job(job_id, progress=30, message="첨부파일 요약 중")

        # 2) 요약
        summarized_by_mail: dict[str, list[dict]] = {}
        for mail_id, items in attachment_texts_by_mail.items():
            summarized_by_mail[mail_id] = [
                {
                    "name": item["name"],
                    "text": _summarize_attachment(item["text"], item["name"])
                }
                for item in items
            ]

        update_job(job_id, progress=50, message="attachment_latest.txt 저장 중")

        # 3) 기록용 attachment_latest.txt 저장
        _write_attachment_file(paths, summarized_by_mail)

        # 기존 본문과 합친 '증분 전용 CSV' 생성
        merged_csv_path = _build_merged_attachment_csv(paths, summarized_by_mail)

        if not merged_csv_path:
            print("[JOB][attachment] 업데이트할 병합 데이터가 없습니다. 종료합니다.")
            update_job(job_id, status="done", message="업데이트할 내용 없음")
            return

        update_job(job_id, progress=60, message="GraphRAG Update 실행 중")

        # 4) graphrag update → json 생성 (마지막 배치일 때만)
        print(f"[JOB][attachment] is_last={is_last}, job_id={job_id}")
        if is_last:
            build_graphrag_update(job_id, paths, env)
            build_graph_json(job_id, paths, env)
        else:
            print(f"[JOB][attachment] 중간 배치 → GraphRAG update 생략, 누적 중")
            _delete_old_update_files(paths)
            mark_attachments_as_processed(paths.GMAIL_ID, attachments)
            update_job(job_id, status="done", message="첨부파일 누적 완료 (중간 배치)")
            return

        # 6) 처리 완료된 이전 update_output 폴더 삭제
        _delete_old_update_files(paths)

        # [추가] 7) 처리 완료된 첨부파일 DB에 기록 (다음 트리거에서 중복 방지)
        mark_attachments_as_processed(paths.GMAIL_ID, attachments)

        update_job(job_id, progress=100, status="done", message="첨부파일 인덱싱 완료")
        print(f"[JOB][attachment] SUCCESS job_id={job_id}")

    except Exception as e:
        error_msg = f"{type(e).__name__}: {e}"
        update_job(job_id, status="failed", message=error_msg)
        print(f"[JOB][attachment][ERROR] job_id={job_id} error={error_msg}")
        traceback.print_exc()

# 기존 mail_latest.csv에서 원본 본문을 읽어와 첨부파일 요약본을 뒤에 붙인 '증분 전용 CSV'를 생성
# 이후 - 새 코드만 남기고 옛날 코드 전부 제거
def _build_merged_attachment_csv(paths, summarized_by_mail: dict[str, list[dict]]):
    # mail_latest.csv + inc_*.csv 전부 읽기 (append 모드에서 새 메일도 포함)
    original_mails = {}
    if not os.path.exists(paths.MAIL_DIR):
        print(f"[AttachmentFile] MAIL_DIR가 없습니다: {paths.MAIL_DIR}")
        return None
    for fname in os.listdir(paths.MAIL_DIR):
        if fname == "mail_latest.csv" or (fname.startswith("inc_") and fname.endswith(".csv") and not fname.startswith("inc_att")):
            csv_path = os.path.join(paths.MAIL_DIR, fname)
            try:
                with open(csv_path, "r", encoding="utf-8") as f:
                    reader = csv.DictReader(f)
                    for row in reader:
                        original_mails[row['id']] = row['text']
            except Exception as e:
                print(f"[AttachmentFile] CSV 읽기 실패: {csv_path} / {e}")

    if not original_mails:
        print(f"[AttachmentFile] 읽을 수 있는 메일 CSV가 없습니다.")
        return None

    csv_rows = []
    for m_id, items in summarized_by_mail.items():
        if m_id in original_mails:
            att_summaries = []
            for item in items:
                att_summaries.append(f"File name: {item['name']}\nSummary: {item['text']}")
            combined_att_text = "\n\n".join(att_summaries)
            combined_text = (
                f"{original_mails[m_id]}\n\n"
                f"[첨부파일 요약]\n"
                f"{combined_att_text}"
            )
            csv_rows.append({"id": m_id, "text": combined_text})
        else:
            print(f"[AttachmentFile] 메일 ID {m_id}를 원본 CSV에서 찾을 수 없습니다.")

    if not csv_rows:
        return None

    new_csv_path = os.path.join(paths.MAIL_DIR, "attachment_latest.csv")
    try:
        # 기존 CSV 읽어서 누적 (중간 배치 내용 보존)
        existing_rows = {}
        if os.path.exists(new_csv_path):
            with open(new_csv_path, "r", encoding="utf-8") as f:
                reader = csv.DictReader(f)
                for row in reader:
                    existing_rows[row["id"]] = row["text"]
        # 새 배치로 갱신 (같은 mail_id면 최신 요약으로 덮어씀)
        for row in csv_rows:
            existing_rows[row["id"]] = row["text"]

        with open(new_csv_path, "w", encoding="utf-8", newline="") as f:
            writer = csv.DictWriter(f, fieldnames=["id", "text"])
            writer.writeheader()
            writer.writerows([{"id": k, "text": v} for k, v in existing_rows.items()])
        print(f"[AttachmentFile] 증분 병합 CSV 생성 완료: {new_csv_path} ({len(existing_rows)}개)")
        return new_csv_path
    except Exception as e:
        print(f"[AttachmentFile] 증분 CSV 생성 중 오류: {e}")
        return None
    
# attachment_latest.txt 저장
def _write_attachment_file(paths, summarized_by_mail: dict[str, list[dict]]):
    att_path = os.path.join(paths.MAIL_DIR, "attachment_latest.txt")

    existing: dict[str, list[dict]] = {}
    if os.path.exists(att_path):
        try:
            with open(att_path, "r", encoding="utf-8") as f:
                raw = f.read()
            existing = _parse_attachment_file(raw)
        except Exception as e:
            print(f"[AttachmentFile] 기존 파일 파싱 실패, 덮어씀: {e}")

    existing.update(summarized_by_mail)

    subjects: dict[str, str] = {}
    if os.path.exists(paths.MAIL_LATEST_PATH):
        with open(paths.MAIL_LATEST_PATH, "r", encoding="utf-8") as f:
            mail_content = f.read()
        for block in mail_content.split(MAIL_BLOCK_SEP):
            id_m = re.search(r"^ID:\s*(.+?)$", block, re.MULTILINE)
            sub_m = re.search(r"제목:\s*(.+?)$", block, re.MULTILINE)
            if id_m and sub_m:
                subjects[id_m.group(1).strip()] = sub_m.group(1).strip()

    lines = []
    for mail_id, items in existing.items():
        for item in items:
            lines.append("[첨부파일 요약]")
            lines.append(f"ID: {mail_id}")
            subject = subjects.get(mail_id, "")
            if subject:
                lines.append(f"제목: {subject}")
            lines.append(f"[File name] {item['name']}")
            lines.append(item['text'])
            lines.append(MAIL_BLOCK_SEP)

    with open(att_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines) + "\n")

    print(f"[AttachmentFile] 저장 완료 → {att_path} ({len(existing)}개 메일)")


# attachment_latest.txt 파싱 → {mail_id: [{name, text}]} 형태로 반환
def _parse_attachment_file(raw: str) -> dict[str, list[dict]]:
    result: dict[str, list[dict]] = {}
    blocks = raw.split(MAIL_BLOCK_SEP)

    for block in blocks:
        block = block.strip()
        if not block:
            continue

        m = re.search(r"^ID:\s*(.+?)$", block, re.MULTILINE)
        if not m:
            continue
        mail_id = m.group(1).strip()

        items = []
        file_blocks = re.split(r"^\[File name\]", block, flags=re.MULTILINE)
        for fb in file_blocks[1:]:
            fb_lines = fb.strip().splitlines()
            if not fb_lines:
                continue
            name = fb_lines[0].strip()
            text = "\n".join(fb_lines[1:]).strip()
            items.append({"name": name, "text": text})

        if items:
            result.setdefault(mail_id, []).extend(items)

    return result

# ============================================================
# [수정] _build_mail_csv: append 모드에서 new_ids가 없을 때 엣지케이스 처리
# 기존: new_ids 없으면 else로 떨어져 mail_latest.csv 전체 덮어씀
# 변경: append 모드에서 new_ids 없으면 None 반환 (CSV 생성 안 함)
#       호출부에서 None 체크 후 graphrag update 생략
# ============================================================
def _build_mail_csv(paths, mode="rewrite", new_ids=None) -> str | None:
    # 1) mail_latest.txt 파싱 → {mail_id: block_text}
    mail_text = _read_latest_text(paths)
    mail_blocks: dict[str, str] = {}

    for block in _split_mail_blocks(mail_text):
        mail_id = _extract_mail_id_from_block(block)
        if mail_id:
            mail_blocks[mail_id] = block.strip()

    # 2) CSV row 생성
    rows = []
    for mail_id, block_text in mail_blocks.items():
        clean_text = block_text.replace(MAIL_BLOCK_SEP, "").strip()
        rows.append({"id": mail_id, "text": clean_text})

    # 3) mode에 따라 저장 대상 결정
    if mode == "append" and new_ids:
        # append + 새 메일 있음: 새 메일만 필터링해서 증분 CSV 생성
        rows = [r for r in rows if r["id"] in new_ids]
        csv_name = f"inc_{datetime.datetime.now().strftime('%Y-%m-%d_%H%M%S')}.csv"

    elif mode == "append" and not new_ids:
        # [수정] append + 새 메일 없음: CSV 생성 불필요 → None 반환
        # 기존에는 else로 떨어져 mail_latest.csv 전체를 덮어쓰는 버그가 있었음
        print("[CSV] append 모드이나 new_ids 없음 → CSV 생성 생략")
        return None

    else:
        # rewrite: 전체를 mail_latest.csv로 저장
        csv_name = "mail_latest.csv"

    csv_path = os.path.join(paths.MAIL_DIR, csv_name)
    with open(csv_path, "w", encoding="utf-8", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=["id", "text"])
        writer.writeheader()
        writer.writerows(rows)

    print(f"[CSV] 생성 완료 → {csv_path} ({len(rows)}개 메일)")
    return csv_path

# 근거메일보기 버튼
def _extract_source_mail_ids(answer: str) -> list:
    return list(set(re.findall(r'ID:\s*([0-9A-Fa-f]{16})', answer)))

# 질의 방법 분류
def _classify_query_method(message: str) -> str:
    prompt = f"""다음 질문이 로컬 검색(특정 메일·인물·날짜·주제)에 적합한지,
                글로벌 검색(전체 경향·요약·패턴·빈도)에 적합한지 판단하라.
                "local" 또는 "global" 중 하나만 반환하라.

                질문: {message}"""

    client = openai.OpenAI(api_key=os.environ.get("GRAPHRAG_API_KEY"))

    res = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=10,
        temperature=0
    )

    method = res.choices[0].message.content.strip().lower()
    print(f"[CLASSIFY] 질의: {message[:30]} → {method}")
    return method if method in ("local", "global") else "local"


# 엔드포인트: POST /extract-calendar
@app.route('/extract-calendar', methods=['POST'])
def extract_calendar():
    data = request.json or {}
    subject = data.get('subject', '')
    body = data.get('body', '')
    result = _convert_to_calendar_json(f"제목: {subject}\n\n{body}")
    return jsonify(result)

# 엔드포인트: POST /run-query-async
@app.route('/run-query-async', methods=['POST'])
def run_query_async():
    data = request.json or {}

    print("[DEBUG] content_type =", request.content_type)
    print("[DEBUG] raw body =", request.data)
    print("[DEBUG] parsed data =", data)

    if data is None:
        return jsonify({'error': 'JSON 본문을 읽지 못했습니다.'}), 400

    message = request.json.get('message', '')
    resMethod = request.json.get('resMethod', 'local')
    resType = request.json.get('resType', 'text')
    gmail_id = data.get('gmail_id', '').strip()

    if not str(message).strip():
        return jsonify({'error': 'message가 비어있습니다.'}), 400

    if not gmail_id:
        return jsonify({'error': 'gmail_id가 비어있습니다.'}), 400

    print("[DEBUG] message =", repr(message))
    print("[DEBUG] gmail_id =", repr(gmail_id))

    job_id = str(uuid.uuid4())[:8]
    create_job(job_id, job_type="query")
    update_job(job_id, status="pending", result=None, resType=resType)


    def _worker():  # 백그라운드 스레드에서 실행되는 실제 작업 함수
        from util.graphrag_query import run_graphrag_query
        try:
            paths = UserPaths(BASE_DIR, gmail_id)
            env = os.environ.copy()
            env["GMAIL_ID"] = gmail_id



            # 날짜 범위 쿼리일 시 parquet 직접 필터링해서 LLM에게 넘기기, 아니면 GraphRAG로 처리
            answer = run_date_range_query(message, paths) # 이게 None이면 GraphRAG로 
            source_ids = []  # 초기화
            if answer is None:
                full_message = message + " 영어 말고 한국어로 답변해줘."

                resMethod = _classify_query_method(message)
                try: # 엔진 객체 직접 호출 방식
                    answer, source_ids = run_graphrag_query(full_message,message, paths, method=resMethod)
                except Exception as e:
                    # API 방식 실패 시 기존 CLI 방식으로 자동 fallback
                    print(f"[ENGINE] API 실패, CLI fallback: {e}")
                    answer = _run_graphrag(full_message,message, resMethod, paths, resType)
                    source_ids = _extract_source_mail_ids(answer)

            if resType.lower() == "calendar":
                result = json.dumps(_convert_to_calendar_json(answer), ensure_ascii=False)
                update_job(job_id, status="done", result=result)
            else:
                result = answer
                update_job(job_id, status="done", result=result, source_ids=source_ids)

        except Exception as e:
            update_job(job_id, status="error", result=str(e))

    threading.Thread(target=_worker, daemon=True).start()
    return jsonify({"jobId": job_id})

# 엔드포인트: GET /job-status/<job_id>
@app.route('/job-status/<job_id>', methods=['GET'])
def job_status(job_id):
    job = get_job(job_id)
    if not job:
        return jsonify({"status": "not_found"}), 404

    if job["status"] == "done" and job["resType"].lower() == "calendar":
        try:
            return jsonify({"status": "done", "data": json.loads(job["result"])})
        except Exception:
            return jsonify({"status": "done", "data": {"events": []}})


    # text 타입: result 필드에 문자열 그대로 반환
    return jsonify({
        "status": job["status"],
        "progress": job.get("progress", 0),
        "message": job.get("message", ""),
        "result": job["result"] or "",
        "source_ids": job.get("source_ids") or [],
    })

# 엔드포인트: GET /indexing-stream (SSE)
# 브라우저가 연결을 유지하면 서버가 인덱싱 progress/완료/실패 이벤트를 즉시 push
# 15초마다 keepalive 전송 (연결 유지용)
@app.route("/indexing-stream", methods=["GET"])
def indexing_stream():
    q = subscribe()

    @stream_with_context
    def generate():
        try:
            while True:
                try:
                    data = q.get(timeout=15)
                    yield f"data: {json.dumps(data, ensure_ascii=False)}\n\n"
                except Exception:
                    yield ": keepalive\n\n"
        finally:
            unsubscribe(q)

    return Response(generate(), content_type="text/event-stream",
                    headers={
                        "Cache-Control": "no-cache",
                        "X-Accel-Buffering": "no",
                        "Connection": "keep-alive",
                        "ngrok-skip-browser-warning": "true",
                    })


# 엔드포인트: GET /indexing-history
@app.route("/indexing-history", methods=["GET"])
def indexing_history():
    """최근 job 상태 목록 반환 (페이지 로드 시 이전 상태 복원용)"""
    all_jobs = get_all_jobs()
    # 최신순 정렬, 최대 20개
    sorted_jobs = sorted(all_jobs.values(), key=lambda j: j.get("created_at", 0), reverse=True)[:20]
    events = []
    for job in sorted_jobs:
        events.append({
            "type": job.get("status", "idle"),
            "job_id": job.get("job_id"),
            "progress": job.get("progress", 0),
            "message": job.get("message", ""),
        })
    return jsonify(events)


# 엔드포인트: POST /run-query (동기 버전)
@app.route('/run-query', methods=['POST'])
def run_query():
    data = request.json or {}
    message = data.get('message', '')
    resMethod = data.get('resMethod', 'local')
    resType = data.get('resType', 'text')
    gmail_id = (data.get('gmail_id') or '').strip().lower()

    print(f'message: {message}')
    print(f'resMethod: {resMethod}')
    print(f'resType: {resType}')

    if not str(message).strip():
        return jsonify({'error': 'message가 비어있습니다.'}), 400
    if not gmail_id:
        return jsonify({'error': 'gmail_id가 비어있습니다.'}), 400

    paths = UserPaths(BASE_DIR, gmail_id)
    message += " 영어 말고 한국어로 답변해줘."

    try:
        answer = _run_graphrag(message, resMethod, paths, resType)
    except RuntimeError as e:
        return jsonify({'error': str(e)}), 500

    if resType.lower() == "calendar":
        return jsonify(_convert_to_calendar_json(answer))

    return jsonify({'result': answer})

# ============================================================
# 엔드포인트: POST /upload
# [수정] 배치 시스템 지원
# - is_last 플래그 수신: 마지막 배치일 때만 GraphRAG 파이프라인 실행
# - 중간 배치: mail_latest.txt에 누적만 하고 GraphRAG 실행 안 함
# - mail_id 기반 중복 블록 체크: rewrite/append 관계없이 항상 적용
# ============================================================
@app.route("/upload", methods=["POST"])
def upload():
    # 1) 데이터 수신
    data = request.json or {}
    filename = data.get("filename") or f"mail_{int(time.time())}.txt"
    content = data.get("content") or ""
    attachments = data.get("attachment") or []
    requested_mode = data.get("syncmode", "append")
    gmail_id = (data.get("gmail_id") or "").strip().lower()
    is_last = data.get("is_last", True)
    batch_offset = data.get("batch_offset", 0)

    paths = UserPaths(BASE_DIR, gmail_id)

    if not str(content).strip():
        return jsonify({"ok": False, "error": "content가 비어있습니다."}), 400
    if not gmail_id:
        return jsonify({"ok": False, "error": "gmail_id가 비어있습니다."}), 400

    print("user gmail id =", gmail_id)
    print(f"[UPLOAD] is_last={is_last}, batch_offset={batch_offset}")

    # append인데 기존 인덱스가 없으면 rewrite로 전환
    fallback_to_rewrite = False
    sync_mode = requested_mode

    if requested_mode == "append" and not _is_index_ready(paths):
        print("[UPLOAD] index not ready -> fallback to rewrite")
        sync_mode = "rewrite"
        fallback_to_rewrite = True

    # 2) 저장 디렉토리 준비
    os.makedirs(paths.MAIL_DIR, exist_ok=True)

    # rewrite 첫 배치(offset=0)에서만 기존 첨부파일 폴더 초기화
    # [수정] 기존: rewrite면 무조건 삭제 → 배치 중간에도 삭제되는 문제
    # 변경: batch_offset=0(첫 배치)일 때만 삭제
    if sync_mode == "rewrite" and batch_offset == 0:
        # rewrite 첫 배치: input 폴더 내 기존 메일 파일 전체 초기화
        # mail_latest.txt, mail_latest.csv, inc_*.txt 등 전부 삭제
        # 이전 데이터가 남아있으면 중복 체크에 걸려 새 메일이 스킵되는 버그 방지
        if os.path.exists(paths.MAIL_DIR):
            for fname in os.listdir(paths.MAIL_DIR):
                fpath = os.path.join(paths.MAIL_DIR, fname)
                try:
                    if os.path.isfile(fpath):  # 파일만 삭제, 폴더는 건너뜀
                        os.remove(fpath)
                except Exception as e:
                    print(f"[CLEAN] 파일 삭제 실패 (무시): {fpath} / {e}")

            print(f"[CLEAN] input 폴더 초기화 완료 (첫 배치): {paths.MAIL_DIR}")
        if os.path.exists(paths.ATTACHMENT_DIR):
            shutil.rmtree(paths.ATTACHMENT_DIR)
            print(f"[CLEAN] attachment 폴더 초기화 완료 (첫 배치): {paths.ATTACHMENT_DIR}")
        # [추가] stats.json 삭제 → 첨부파일 트리거가 인덱스 없음으로 판단해 거절됨
        # rewrite 완료 전에 첨부파일이 먼저 처리되는 문제 방지
        stats_path = os.path.join(paths.GRAPHRAG_ROOT, "output", "stats.json")
        if os.path.exists(stats_path):
            try:
                os.remove(stats_path)
                print(f"[CLEAN] stats.json 삭제 완료 (rewrite 시작)")
            except Exception as e:
                print(f"[CLEAN] stats.json 삭제 실패 (무시): {e}")

        try:
            from util.database.db_writer import get_latest_user_record
            latest_user = get_latest_user_record(gmail_id)
            if latest_user:
                conn = get_db_connection()
                cursor = conn.cursor()
                cursor.execute(
                    "DELETE FROM processed_attachments WHERE user_account_id = %s AND update_date = %s",
                    (latest_user["user_account_id"], latest_user["update_date"])
                )
                conn.commit()
                cursor.close()
                conn.close()
            print(f"[CLEAN] processed_attachments DB 초기화 완료 (gmail_id={gmail_id})")
        except Exception as e:
            print(f"[CLEAN] processed_attachments DB 초기화 실패 (무시): {e}")

    # 3) 원본 메일 텍스트 저장
    # [수정] rewrite 모드 배치 누적 버그 수정
    # 기존: rewrite 모드에서 배치마다 "mail_latest.txt"를 "w" 모드로 열어 덮어씀
    #       → 배치2가 오면 배치1 내용이 사라지고 배치2만 남는 문제
    # 변경: rewrite 중간 배치는 "mail_latest.txt"에 "a" 모드로 이어붙임
    #       첫 배치(batch_offset=0)일 때만 파일을 비우고 시작
    #       append 모드는 기존 방식 유지 (inc_*.txt로 별도 저장)
    if sync_mode != "rewrite":  # rewrite는 여기서 파일 안 씀
        file_path = os.path.join(paths.MAIL_DIR, filename)
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(content)

    extracted_count = 0
    failed_attachments = []

    # 4) 첨부파일 메타데이터 카운트 (원본은 트리거가 별도 전송)
    for file_info in attachments:
        f_name = file_info.get("name") or "attachment.bin"
        mail_id = str(file_info.get("mail_id") or "").strip()
        if f_name and mail_id:
            extracted_count += 1

    # 5) 로그
    print(f"[UPLOAD] Received filename: {filename}")
    print(f"[UPLOAD] Content length: {len(content)}")
    print(f"[UPLOAD] Attachment count received: {len(attachments)}")
    print(f"[UPLOAD] Attachment extracted count: {extracted_count}")
    print(f"[UPLOAD] Requested mode: {requested_mode}")
    print(f"[UPLOAD] Actual mode: {sync_mode}")
    print(f"[UPLOAD] is_last: {is_last}")
    print("[UPLOAD] cwd:", os.getcwd())

    added_count  = 0
    skipped_count = 0
    saved_mail_path = ""

    # ============================================================
    # [수정] 메일 텍스트 누적 로직
    # rewrite: 파일에 직접 이어붙이므로 중복 체크만 수행 (mail_latest.txt 재작성 불필요)
    # append: 기존 방식 유지 (existing_text 읽어서 합친 후 mail_latest.txt 저장)
    # 공통: mail_id 기반 중복 체크 → 배치 재시도 시 중복 삽입 방지
    # ============================================================

    # 기존 mail_latest.txt에서 이미 저장된 mail_id 추출 (중복 방지용)
    # rewrite 첫 배치: 파일을 새로 쓰는 시점이므로 기존 내용 무시
    # → 기존 파일의 mail_id를 읽으면 전부 중복으로 판단해서 스킵되는 버그 방지
    if sync_mode == "rewrite" and batch_offset == 0:
        existing_text = ""
        existing_ids  = set()
    else:
        existing_text = _read_latest_text(paths)
        existing_ids  = _extract_message_ids(existing_text)

    new_blocks    = _split_mail_blocks(content)
    append_blocks = []

    for block in new_blocks:
        msg_id = _extract_mail_id_from_block(block)
        if not msg_id:
            skipped_count += 1
            continue
        if msg_id in existing_ids:
            skipped_count += 1
            continue
        append_blocks.append(block.strip())
        existing_ids.add(msg_id)

    added_count = len(append_blocks)

    # rewrite 첫 배치: 증분 파일 초기화
    if sync_mode == "rewrite" and batch_offset == 0:
        _delete_incremental_files(paths)

    if batch_offset == 0:  # rewrite/append 공통으로 밖으로 꺼냄
        batch_job_id = "batch_" + gmail_id
        create_job(batch_job_id, job_type="batch")
        update_job(batch_job_id, status="running", message="배치 진행 중")
        print(f"[UPLOAD] 배치 시작 job 생성: {batch_job_id}")

    if append_blocks:
        append_blocks.sort(key=_extract_block_for_sort, reverse=True)
        inc_content = "\n\n".join(append_blocks).strip() + "\n"
        # [수정] rewrite: 파일에 직접 이어붙이는 방식으로 변경
        # 파일 저장은 위(3번)에서 이미 완료됨 ("a" 모드로 이어붙임)
        # 여기서는 _renumber_mail_blocks만 적용해서 최종 정리
        # 단, 마지막 배치일 때만 번호 재정렬 (중간 배치는 불완전한 상태)
        if sync_mode == "rewrite":
            with open(paths.MAIL_LATEST_PATH, "a", encoding="utf-8") as f:
                f.write(inc_content)  # 정제된 블록만 이어붙임
            if is_last:
                final_text = _read_latest_text(paths)
                all_blocks = _split_mail_blocks(final_text)
                all_blocks.sort(key=_extract_block_for_sort, reverse=True)  # 날짜 정렬
                sorted_text = "\n\n".join(b.strip() for b in all_blocks).strip() + "\n"
                with open(paths.MAIL_LATEST_PATH, "w", encoding="utf-8") as f:
                    f.write(_renumber_mail_blocks(sorted_text))
        else:
            # append: 기존 내용 앞에 새 메일 추가 후 mail_latest.txt 저장
            existing_lines = existing_text.splitlines()
            existing_clean = "\n".join(existing_lines).lstrip("\n")
            updated_content = inc_content + "\n" + existing_clean
            with open(paths.MAIL_LATEST_PATH, "w", encoding="utf-8") as f:
                f.write(_renumber_mail_blocks(updated_content.strip()))

        saved_mail_path = paths.MAIL_LATEST_PATH

        # 새로 추가된 메일 ID 수집
        new_ids = set()
        for block in append_blocks:
            mid = _extract_mail_id_from_block(block)
            if mid:
                new_ids.add(mid)

        # statics 파이프라인
        statics_job_id = str(uuid.uuid4())[:8]
        create_job(statics_job_id, job_type="statics")
        
        if is_last and sync_mode == "rewrite":
            final_text = _read_latest_text(paths)
            statics_blocks = _split_mail_blocks(final_text)
            statics_blocks = [b for b in statics_blocks if _extract_mail_id_from_block(b)]
        else:
            statics_blocks = append_blocks

        start_statics_pipeline_background(
            statics_job_id, paths,
            mode="rewrite" if sync_mode == "rewrite" else "append"
        )

    else:
        saved_mail_path = ""
        new_ids = set()

    print("[UPLOAD] added:", added_count)
    print("[UPLOAD] skipped:", skipped_count)
    if saved_mail_path:
        print("[UPLOAD] saved mail path:", os.path.abspath(saved_mail_path))

    # ============================================================
    # [수정] GraphRAG 파이프라인 실행 조건
    # 기존: 업로드 때마다 GraphRAG 실행
    # 변경: is_last=True 일 때만 실행
    #       - 중간 배치(is_last=False): mail_latest.txt 누적만, GraphRAG 실행 안 함
    #       - 마지막 배치(is_last=True): 전체 누적 텍스트로 GraphRAG 실행
    #       - 배치 시스템 없는 기존 단일 호출: is_last 기본값 True → 기존 동작 유지
    # ============================================================
    graph_job_id = str(uuid.uuid4())[:8]

    if not is_last:
        # 중간 배치: GraphRAG 실행 안 함, 누적만
        print(f"[UPLOAD] 중간 배치 (is_last=False) → GraphRAG 실행 생략, 누적 중")
        return jsonify({
            "ok": True,
            "requested_mode": requested_mode,
            "actual_mode": sync_mode,
            "is_last": is_last,
            "fallback_to_rewrite": fallback_to_rewrite,
            "added_count": added_count,
            "skipped_count": skipped_count,
            "attachment_received_count": len(attachments),
            "attachment_extracted_count": extracted_count,
        })

    # 마지막 배치: GraphRAG 파이프라인 실행
    batch_job_id = "batch_" + gmail_id
    update_job(batch_job_id, status="done", message="배치 완료")
    print(f"[UPLOAD] 배치 완료 job 닫기: {batch_job_id}")

    # 마지막 배치: GraphRAG 파이프라인 실행
    if sync_mode == "rewrite":
        create_job(graph_job_id, job_type="index")
        update_job(graph_job_id, message="업로드 완료, 그래프 파이프라인 시작")
    else:
        create_job(graph_job_id, job_type="update")
        update_job(graph_job_id, message="업로드 완료, 그래프 업데이트 파이프라인 시작")

    env = os.environ.copy()
    env["PYTHONUNBUFFERED"] = "1"

    if sync_mode == "rewrite":
        update_dir = os.path.join(paths.GRAPHRAG_ROOT, "update_output")
        if os.path.exists(update_dir):
            shutil.rmtree(update_dir)
            print(f"[CLEAN] update_output 삭제 완료: {update_dir}")

        # [순서 주석] _build_mail_csv는 동기 실행 후 GraphRAG 스레드 시작
        # → CSV 파일이 완전히 쓰인 뒤 GraphRAG가 읽도록 순서 보장
        _build_mail_csv(paths)
        # rewrite 배치 완료 시 총 누적 메일 수로 기록 (마지막 배치 added_count만 넘기면 일부만 저장되는 버그 방지)
        final_text = _read_latest_text(paths)
        total_mail_count = len([b for b in _split_mail_blocks(final_text) if _extract_mail_id_from_block(b)])
        start_graph_pipeline_background(graph_job_id, paths, env, added_count=total_mail_count)

    else:  # append
        if new_ids:
            # [수정] _build_mail_csv 반환값 None 체크 추가
            # new_ids 없을 때 None 반환하도록 수정했으므로 None이면 update 생략
            csv_path = _build_mail_csv(paths, mode="append", new_ids=new_ids)
            if csv_path:
                start_graph_update_pipeline_background(graph_job_id, paths, env)
            else:
                update_job(graph_job_id, status="done", message="CSV 없음, 업데이트 생략")
                print("[UPLOAD] CSV 생성 실패 → graphrag update 생략")
        else:
            update_job(graph_job_id, status="done", message="추가된 새 메일 없음, 업데이트 생략")
            print("[UPLOAD] new_ids 없음 → graphrag update 생략")

    return jsonify({
        "ok": True,
        "requested_mode": requested_mode,
        "job_id": graph_job_id,
        "actual_mode": sync_mode,
        "fallback_to_rewrite": fallback_to_rewrite,
        "is_last": is_last,
        "latest_path": os.path.abspath(paths.MAIL_LATEST_PATH),
        "saved_mail_path": os.path.abspath(saved_mail_path) if saved_mail_path else "",
        "attachment_dir": os.path.abspath(paths.ATTACHMENT_DIR),
        "content_length": len(content),
        "added_count": added_count,
        "skipped_count": skipped_count,
        "attachment_received_count": len(attachments),
        "attachment_extracted_count": extracted_count,
        "failed_attachments": failed_attachments,
    })

# 엔드포인트: GET /graph-data
@app.route("/graph-data", methods=["GET", "OPTIONS"])
def graph_data():
    if request.method == "OPTIONS":
        return "", 200

    gmail_id = (request.args.get("gmail_id") or "").strip().lower()

    if not gmail_id:
        return jsonify({"ok": False, "error": "gmail_id가 비어있습니다."}), 400

    paths = UserPaths(BASE_DIR, gmail_id)

    if not os.path.exists(paths.GRAPH_JSON_PATH):
        return jsonify({"nodes": [], "edges": [], "error": "graph json not found"}), 200

    try:
        with open(paths.GRAPH_JSON_PATH, "rb") as f:
            raw = f.read().rstrip(b'\x00')  # null 바이트 제거 (비정상 종료 방어)
        data = json.loads(raw.decode("utf-8"))
        print(f"[GRAPH-DATA] 반환: {len(data.get('nodes', []))} 노드")
        return jsonify(data)
    except Exception as e:
        print(f"[GRAPH-DATA] 에러: {e}")
        return jsonify({"nodes": [], "edges": [], "error": str(e)}), 500

# 엔드포인트: GET /graph-view
@app.route("/graph-view", methods=["GET"])
def graph_view():
    return send_from_directory(
        os.path.join(os.path.dirname(__file__), "json"),
        "graph_view.html"
    )

# 공유 그래프 렌더링 함수
@app.route('/graph-render.js')
def graph_render_js():
    return send_from_directory(
        os.path.join(os.path.dirname(__file__), "json"),
        "graph-render.js"
    )

# 엔드포인트: GET /index-status
@app.route("/index-status", methods=["GET"])
def index_status():
    gmail_id = (request.args.get("gmail_id") or "").strip().lower()
    if not gmail_id:
        return jsonify({"error": "gmail_id가 비어있습니다."}), 400
    paths = UserPaths(BASE_DIR, gmail_id)
    return jsonify({"indexed": _is_index_ready(paths)})

# 엔드포인트: GET /init  — localStorage에 flask_url 자동 저장 후 대시보드로 이동
@app.route('/init')
def init_storage():
    from flask import request as _req
    origin = _req.host_url.rstrip('/')
    return f"""<!DOCTYPE html><html><head><meta charset="utf-8"><title>Initializing...</title></head>
<body>
<script>
  localStorage.setItem('gw_flask_url', {repr(origin)});
  window.location.replace('/dashboard/');
</script>
<p>설정 중... 자동으로 이동합니다.</p>
</body></html>""", 200, {{'Content-Type': 'text/html; charset=utf-8'}}

# 엔드포인트: GET /dashboard/
@app.route('/dashboard/', defaults={'path': 'production/index.html'})
@app.route('/dashboard/<path:path>')
def dashboard(path):
    dist_dir = os.path.join(os.path.dirname(__file__), 'web', 'dist')
    if not path.startswith('production/') and path.endswith('.html'):
        path = 'production/' + path
    return send_from_directory(dist_dir, path)

@app.route('/assets/<path:path>')
def static_assets(path):
    dist_dir = os.path.join(os.path.dirname(__file__), 'web', 'dist', 'assets')
    return send_from_directory(dist_dir, path)

@app.route('/js/<path:path>')
def static_js(path):
    dist_dir = os.path.join(os.path.dirname(__file__), 'web', 'dist', 'js')
    return send_from_directory(dist_dir, path)

@app.route('/fonts/<path:path>')
def static_fonts(path):
    dist_dir = os.path.join(os.path.dirname(__file__), 'web', 'dist', 'fonts')
    return send_from_directory(dist_dir, path)

# 엔드포인트: POST /calendar-events
@app.route('/calendar-events', methods=['POST'])
def calendar_events():
    data = request.json or {}
    res = requests.post(WEBAPP_URL, json=data, allow_redirects=True)
    print("[calendar] status:", res.status_code)
    print("[calendar] response:", res.text[:500])
    try:
        return jsonify(res.json())
    except Exception:
        return jsonify({"events": [], "error": res.text[:200]}), 200

# 엔드포인트: POST /labels-proxy
@app.route('/labels-proxy', methods=['POST'])
def labels_proxy():
    data = request.json or {}
    print("[labels] 받은 data:", data)
    try:
        res = requests.post(WEBAPP_URL, json=data, allow_redirects=False, timeout=30)
        print("[labels] 1차 status:", res.status_code)

        if res.status_code in (301, 302, 303, 307, 308):
            location = res.headers.get("Location")
            print("[labels] redirect →", location)
            if res.status_code in (307, 308):
                res = requests.post(location, json=data, allow_redirects=True, timeout=30)
            else:
                res = requests.get(location, allow_redirects=True, timeout=30)
            print("[labels] 2차 status:", res.status_code)

        content_type = res.headers.get("Content-Type", "")
        if "text/html" in content_type:
            msg = re.search(r'class="errorMessage"[^>]*>(.*?)</div>', res.text, re.DOTALL)
            if msg:
                error_text = re.sub(r'<[^>]+>', '', msg.group(1)).strip()
            else:
                plain = re.sub(r'<[^>]+>', ' ', res.text)
                plain = re.sub(r'\s+', ' ', plain).strip()
                error_text = plain[:200] if plain else 'Apps Script 연결 오류'
            print("[labels] GAS 오류 메시지:", error_text)
            return jsonify({"ok": False, "error": error_text}), 200

        result = res.json()
        print("[labels] GAS 응답:", result)
        return jsonify(result)

    except Exception as e:
        return jsonify({"ok": False, "error": str(e)}), 500


# 라벨 서치 프롬포트
with open(os.path.join("parquet_template", "prompts", "label_search_prompt.txt"), "r", encoding="utf-8") as _f:
    LABEL_SEARCH_PROMPT = _f.read().strip()

# 엔드포인트: POST /label-route
@app.route("/label-route", methods=["POST"])
def label_route_deprecated():
    data = request.json or {}
    with app.test_request_context(
        '/label-query', method='POST',
        json=data, content_type='application/json'
    ):
        result = label_query()
    result_data = result.get_json()
    if not result_data.get("ok"):
        return jsonify(result_data)
    return jsonify({"ok": True, "intent": result_data.get("intent", "query")})

# 엔드포인트: POST /label-query (라벨 특화 질의 - 의도 분류 + OpenAI Function Calling)
@app.route("/label-query", methods=["POST"])
def label_query():
    data        = request.json or {}
    user_input  = data.get("userInput", "").strip()
    label_names = data.get("labels", [])

    if not user_input:
        return jsonify({"ok": False, "error": "userInput이 비어있습니다."}), 400

    system_content = (
        "사용자의 Gmail 관련 요청을 분석하세요.\n\n"

            "【작업 요청】메일을 검색·분류·삭제·이동·라벨 추가·라벨 제거하는 요청이면 "
        "반드시 아래 함수 중 하나를 호출하세요.\n"
        "  - search_emails : 메일 검색 후 라벨 적용\n"
        "  - apply_label   : 이미 선택된 메일에 라벨만 적용\n"
        "  - trash_emails  : 메일을 휴지통으로 이동\n"
        "  - remove_label  : 메일에서 라벨 제거\n"
        "  - create_label  : 새 라벨 생성\n\n"

        "【정보 질문】메일 내용이나 현황에 대한 질문 "
        "(예: '받은 메일 중 중요한 게 뭐야?', '최근 메일 요약해줘', "
        "'어떤 라벨에 메일이 많아?', '내 메일함 분석해줘') 이면 "
        "함수를 호출하지 마세요. 아무 함수도 호출하지 않으면 "
        "시스템이 자동으로 GraphRAG 지식 베이스 검색으로 답변합니다.\n\n"
        "판단이 애매할 때는 작업 요청으로 처리하세요."
    )

    if label_names:
        system_content += f"\n\n현재 사용자의 라벨 목록: {', '.join(label_names)}"

    tools = [
        {
            "type": "function",
            "function": {
                "name": "search_emails",
                "description": "사용자 요청에서 Gmail 검색 키워드와 적용할 라벨명을 추출합니다.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "query": {"type": "string", "description": "Gmail 검색에 사용할 키워드"},
                        "label_to_apply": {"type": "string", "description": "검색된 메일에 적용할 라벨명. '현재 사용자의 라벨 목록'에 있는 이름을 그대로 사용하세요. 목록에 없거나 언급이 없으면 빈 문자열."}
                    },
                    "required": ["query", "label_to_apply"]
                }
            }
        },
        {
            "type": "function",
            "function": {
                "name": "apply_label",
                "description": "선택된 메일 ID 목록에 라벨을 적용합니다.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "message_ids": {"type": "array", "items": {"type": "string"}, "description": "라벨을 적용할 메일 ID 목록"},
                        "label_name": {"type": "string", "description": "적용할 라벨명"}
                    },
                    "required": ["message_ids", "label_name"]
                }
            }
        },
        {
            "type": "function",
            "function": {
                "name": "trash_emails",
                "description": "삭제하거나 휴지통으로 이동할 메일의 검색 키워드를 추출합니다.",
                "parameters": {
                    "type": "object",
                    "properties": {

                        "query": {
                            "type": "string",
                            "description": "삭제할 메일을 찾기 위한 Gmail 검색 키워드"
                        }
                    },
                    "required": ["query"]
                }
            }
        },
        {
            "type": "function",
            "function": {
                "name": "remove_label",

                "description": "메일에서 특정 라벨을 제거합니다. '라벨 빼줘', '라벨 제거해줘' 같은 요청에 사용합니다.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "query": {
                            "type": "string",
                            "description": "라벨을 제거할 메일을 찾기 위한 Gmail 검색 키워드"
                        },
                        "label_name": {
                            "type": "string",
                            "description": "제거할 라벨명. 언급이 없으면 빈 문자열."
                        }
                    },
                    "required": ["query", "label_name"]
                }
            }
        },
        {
            "type": "function",
            "function": {
                "name": "create_label",
                "description": "새 Gmail 라벨을 생성합니다. '라벨 만들어줘', '라벨 추가해줘', '~라벨 생성해줘' 같은 요청에 사용합니다.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "label_name": {
                            "type": "string",
                            "description": "생성할 라벨명 (예: 산학공동연구)"
                        }
                    },
                    "required": ["label_name"]
                }
            }
        }
    ]

    try:
        client = openai.OpenAI(api_key=os.environ.get("GRAPHRAG_API_KEY"))
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": system_content},
                {"role": "user",   "content": user_input}
            ],
            tools=tools,
            tool_choice="auto",
            temperature=0
        )
        tool_calls = response.choices[0].message.tool_calls

        if not tool_calls:
            print("[label-query] FC 없음 → intent=query")
            return jsonify({"ok": True, "intent": "query"})
        
        if len(tool_calls) == 1:
            action = tool_calls[0].function.name
            params = json.loads(tool_calls[0].function.arguments)
            print(f"[label-query] 단일 액션: {action} / params: {params}")
            return jsonify({
                "ok":     True,
                "intent": "action",
                "action": action,
                "params": params
            })
        
        actions = [
            {"action": tc.function.name, "params": json.loads(tc.function.arguments)}
            for tc in tool_calls
        ]
        print(f"[label-query] 복합 액션: {[a['action'] for a in actions]}")
        return jsonify({
            "ok":      True,
            "intent":  "action",
            "action":  "multi",
            "actions": actions
        })

    except Exception as e:
        return jsonify({"ok": False, "error": str(e)}), 500


# ============================================================
# 엔드포인트: POST /upload-attachments
# [수정] 중복 처리 방지 로직 추가
# 기존: 10분마다 전체 첨부파일을 무조건 처리
# 변경: DB 조회로 이미 처리된 (gmail_id, mail_id, filename) 조합 필터링 후 처리
#       처리 완료 후 DB에 기록 → 다음 트리거에서 중복 처리 방지
# ============================================================
@app.route("/upload-attachments", methods=["POST"])
def upload_attachments():
    # 1) 데이터 수신
    data = request.json or {}
    gmail_id = (data.get("gmail_id") or "").strip().lower()
    attachments = data.get("attachments") or []

    if not gmail_id:
        return jsonify({"ok": False, "error": "gmail_id가 비어있습니다."}), 400
    
    if not attachments:
        # attachments 없이 is_last=true만 온 경우 → GraphRAG update 트리거
        is_last = data.get("is_last", False)
        if is_last:
            # 이미 누적된 attachment_latest.csv로 GraphRAG update 실행
            paths = UserPaths(BASE_DIR, gmail_id)
            if os.path.exists(os.path.join(paths.MAIL_DIR, "attachment_latest.csv")):
                job_id = str(uuid.uuid4())[:8]
                create_job(job_id, job_type="attachment")
                env = os.environ.copy()
                env["PYTHONUNBUFFERED"] = "1"
                from util.jobs.job_run import build_graphrag_update, build_graph_json
                def _finish():
                    build_graphrag_update(job_id, paths, env)
                    build_graph_json(job_id, paths, env)
                    _delete_old_update_files(paths)
                    update_job(job_id, status="done", message="첨부파일 인덱싱 완료")
                    print(f"[JOB][attachment] SUCCESS job_id={job_id}")
                threading.Thread(target=_finish, daemon=True).start()
                return jsonify({"ok": True, "message": "finish signal received"})
        return jsonify({"ok": False, "error": "attachments가 비어있습니다."}), 400
    
    paths = UserPaths(BASE_DIR, gmail_id)

    # 2) 메일 인덱스가 준비되지 않았으면 거절
    # 메일 본문 인덱싱 완료 전에 첨부파일 처리하면 불완전한 그래프에 update가 붙는 문제 방지
    # 10분 트리거가 다음번에 재시도함
    if not _is_index_ready(paths):
        print(f"[upload-attachments] 메일 인덱스 미준비 → 요청 거절, 다음 트리거에서 재시도")
        return jsonify({"ok": False, "error": "메일 인덱스 미준비, 다음 트리거에서 재시도됩니다."}), 409

    # 3) 인덱싱/업데이트 중이면 거절 (graphrag 동시 실행 방지)
    running_jobs = [j for j in get_all_jobs().values()
                if j.get("status") == "running"
                and j.get("job_type") in ("index", "update", "batch")]
    
    if running_jobs:
        print(f"[upload-attachments] 인덱싱 진행 중 → 요청 거절, 다음 트리거에서 재시도")
        return jsonify({"ok": False, "error": "인덱싱 진행 중, 다음 트리거에서 재시도됩니다."}), 409

    # [추가] 4) 이미 처리된 첨부파일 필터링
    is_last = data.get("is_last", True)
    unprocessed = filter_unprocessed_attachments(gmail_id, attachments)

    if not unprocessed:
        print(f"[upload-attachments] 모두 이미 처리된 첨부파일 → 스킵")
        if is_last and os.path.exists(os.path.join(paths.MAIL_DIR, "attachment_latest.csv")):
            job_id = str(uuid.uuid4())[:8]
            create_job(job_id, job_type="attachment")
            env = os.environ.copy()
            env["PYTHONUNBUFFERED"] = "1"
            from util.jobs.job_run import build_graphrag_update, build_graph_json
            def _finish():
                build_graphrag_update(job_id, paths, env)
                build_graph_json(job_id, paths, env)
                _delete_old_update_files(paths)
                update_job(job_id, status="done", message="첨부파일 인덱싱 완료")
                print(f"[JOB][attachment] SUCCESS job_id={job_id}")
            threading.Thread(target=_finish, daemon=True).start()
            return jsonify({"ok": True, "message": "모두 처리됨, finish 실행"})
        return jsonify({"ok": True, "skipped": len(attachments), "message": "모두 이미 처리된 첨부파일"})

    # 4) 즉시 200 응답 (Apps Script 타임아웃 방지)
    job_id = str(uuid.uuid4())[:8]
    create_job(job_id, job_type="attachment")
    update_job(job_id, message="첨부파일 수신 완료, 백그라운드 처리 시작")

    env = os.environ.copy()
    env["PYTHONUNBUFFERED"] = "1"

    # 5) 백그라운드에서 처리 (미처리 첨부파일만 전달)
    t = threading.Thread(
        target=_run_attachment_pipeline,
        args=(job_id, paths, unprocessed, env, is_last),
        daemon=True
    )
    t.start()

    return jsonify({
        "ok": True,
        "job_id": job_id,
        "attachment_count": len(unprocessed),
        "skipped_count": len(attachments) - len(unprocessed),
    })


# 웹앱용 통계 라우트
@app.route("/mail-stats", methods=["POST"])
def send_mail_stats():
    data = request.json or {}
    gmail_id = data.get("gmail_id", "").strip()
    if not gmail_id:
        return jsonify({"error": "gmail_id is required"}), 400
    paths = UserPaths(BASE_DIR, gmail_id)
    print(f"[MAIL_STATS] gmail_id={gmail_id}")
    print(f"[MAIL_STATS] path={paths.USER_ROOT}")
    return jsonify({"gmail_id": gmail_id, "data": get_mail_stats(paths)})

@app.route("/mail-date-range", methods=["POST"])
def send_mail_date_range():
    data = request.json or {}
    gmail_id = data.get("gmail_id", "").strip()
    if not gmail_id:
        return jsonify({"error": "gmail_id is required"}), 400
    return jsonify({"gmail_id": gmail_id, "data": get_mail_date_range(gmail_id)})

@app.route("/keyword-stats", methods=["POST"])
def send_keyword_stats():
    data = request.json or {}
    gmail_id = data.get("gmail_id", "").strip()
    if not gmail_id:
        return jsonify({"error": "gmail_id is required"}), 400
    paths = UserPaths(BASE_DIR, gmail_id)
    return jsonify({"gmail_id": gmail_id, "data": get_keyword_stats(paths)})

@app.route("/keyword-by-person-date", methods=["POST"]) # 각 사람마다 주고받은 메일의 키위드 리턴
def keyword_by_person_date():
    data = request.json or {}
    gmail_id = data.get("gmail_id", "").strip()
    person_gmail_id = data.get("person_gmail_id", "").strip()
    # 시간 범위 내에 있는 메일의 키워드들을 추출
    start_date = data.get("start_date", "").strip()
    end_date = data.get("end_date", "").strip()

    if not gmail_id:
        return jsonify({"error": "gmail_id is required"}), 400
    if not person_gmail_id:
        return jsonify({"error": "person_gmail_id is required"}), 400
    if not start_date or not end_date:
        return jsonify({"error": "start_date and end_date are required"}), 400

    try:
        keywords = get_keywords_by_person_date(gmail_id, person_gmail_id, start_date, end_date)
        return jsonify({"keywords": keywords})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/high_affinity_person_stats", methods=["POST"])
def send_high_affinity_person_stats():
    data = request.json or {}
    gmail_id = data.get("gmail_id", "").strip()
    if not gmail_id:
        return jsonify({"error": "gmail_id is required"}), 400
    paths = UserPaths(BASE_DIR, gmail_id)
    return jsonify({"gmail_id": gmail_id, "data": get_high_affinity_person_stats(paths)})

@app.route("/user_rating_stats", methods=["POST"])
def send_user_rating_stats():
    data = request.json or {}
    gmail_id = data.get("gmail_id", "").strip()
    if not gmail_id:
        return jsonify({"error": "gmail_id is required"}), 400
    paths = UserPaths(BASE_DIR, gmail_id)
    return jsonify({"gmail_id": gmail_id, "data": get_user_rating_stats()})

@app.route("/mail_sync_stats", methods=["POST"])
def send_mail_sync_stats():
    data = request.json or {}
    gmail_id = data.get("gmail_id", "").strip()
    if not gmail_id:
        return jsonify({"error": "gmail_id is required"}), 400
    paths = UserPaths(BASE_DIR, gmail_id)
    return jsonify({"gmail_id": gmail_id, "data": get_mail_sync_stats(paths)})

@app.route("/mail-exchange-stats", methods=["POST"])
def send_mail_exchange_stats():
    data = request.json or {}
    gmail_id       = data.get("gmail_id", "").strip()
    person_mail_id = data.get("person_gmail_id", "").strip()
    start_date     = data.get("start_date", "").strip()
    end_date       = data.get("end_date", "").strip()

    if not gmail_id:
        return jsonify({"error": "gmail_id is required"}), 400
    if not person_mail_id:
        return jsonify({"error": "person_gmail_id is required"}), 400
    if not start_date or not end_date:
        return jsonify({"error": "start_date and end_date are required"}), 400

    return jsonify({"data": get_mail_exchange_stats(gmail_id, person_mail_id, start_date, end_date)})

@app.route("/mail-person-sent-stats", methods=["POST"])
def send_mail_person_sent_stats():
    data = request.json or {}
    gmail_id   = data.get("gmail_id", "").strip()
    start_date = data.get("start_date", "").strip()
    end_date   = data.get("end_date", "").strip()

    if not gmail_id:
        return jsonify({"error": "gmail_id is required"}), 400
    if not start_date or not end_date:
        return jsonify({"error": "start_date and end_date are required"}), 400

    return jsonify({"gmail_id": gmail_id, "data": get_date_range_person_stats(gmail_id, start_date, end_date, "sent")})

@app.route("/mail-person-received-stats", methods=["POST"])
def send_mail_person_received_stats():
    data = request.json or {}
    gmail_id   = data.get("gmail_id", "").strip()
    start_date = data.get("start_date", "").strip()
    end_date   = data.get("end_date", "").strip()

    if not gmail_id:
        return jsonify({"error": "gmail_id is required"}), 400
    if not start_date or not end_date:
        return jsonify({"error": "start_date and end_date are required"}), 400

    return jsonify({"gmail_id": gmail_id, "data": get_date_range_person_stats(gmail_id, start_date, end_date, "received")})

@app.route("/mail-summaries", methods=["POST"])
def send_mail_summaries():
    data = request.json or {}
    gmail_id     = data.get("gmail_id", "").strip()
    summary_type = data.get("type", "").strip()

    if not gmail_id:
        return jsonify({"error": "gmail_id is required"}), 400
    if summary_type not in ("monthly", "yearly"):
        return jsonify({"error": "type must be 'monthly' or 'yearly'"}), 400

    paths = UserPaths(BASE_DIR, gmail_id)
    if not os.path.exists(paths.MAIL_SUMMARIES_PATH):
        return jsonify({"error": "summaries not generated yet"}), 404

    with open(paths.MAIL_SUMMARIES_PATH, "r", encoding="utf-8") as f:
        summaries = json.load(f)

    return jsonify({summary_type: summaries.get(summary_type, {})})

# 연락처 프록시
@app.route('/contacts-proxy', methods=['POST'])
def contacts_proxy():
    data = request.get_json() or {}
    action = data.get('action', '')
    gmail_id = (data.get('gmail_id') or '').strip().lower()

    if not gmail_id:
        return jsonify({'ok': False, 'error': 'gmail_id가 비어있습니다.'}), 400

    paths = UserPaths(BASE_DIR, gmail_id)

    if action == 'getFrequentContacts':
        max_results = int(data.get('maxResults', 100))
        try:
            if not os.path.exists(paths.MAIL_CONTACTS_PATH):
                return jsonify({'ok': True, 'contacts': []})
            with open(paths.MAIL_CONTACTS_PATH, 'r', encoding='utf-8') as f:
                stats = json.load(f)
            result = []
            for email, info in stats.items():
                count = info.get('sent', 0) + info.get('received', 0)
                result.append({
                    'email': email,
                    'name': info.get('name', '') or email.split('@')[0],
                    'count': count,
                    'lastMailAt': None,
                })
            result.sort(key=lambda x: -x['count'])
            return jsonify({'ok': True, 'contacts': result[:max_results]})
        except Exception as e:
            return jsonify({'ok': False, 'error': str(e)})

    elif action == 'getMailHistory':
        email = (data.get('email') or '').strip()
        if not email:
            return jsonify({'ok': False, 'error': 'email이 비어있습니다.'}), 400
        try:
            if not os.path.exists(paths.MAIL_CONTACTS_PATH):
                return jsonify({'ok': True, 'sentCount': 0, 'receivedCount': 0})
            with open(paths.MAIL_CONTACTS_PATH, 'r', encoding='utf-8') as f:
                stats = json.load(f)
            info = stats.get(email, {})
            return jsonify({
                'ok': True,
                'sentCount': info.get('sent', 0),
                'receivedCount': info.get('received', 0),
            })
        except Exception as e:
            return jsonify({'ok': False, 'error': str(e)})

    return jsonify({'ok': False, 'error': f'unknown action: {action}'})

# 메일 보내기
@app.route('/send-mail', methods=['POST', 'OPTIONS'])
def send_mail():
    if request.method == 'OPTIONS':
        return '', 204
    data = request.get_json() or {}
    try:
        res = requests.post(WEBAPP_URL, json={
            'action':  'sendMail',
            'to':      data.get('to'),
            'subject': data.get('subject'),
            'body':    data.get('body'),
        }, allow_redirects=False, timeout=30)
        if res.status_code in (301, 302, 303, 307, 308):
            location = res.headers.get('Location')
            res = requests.get(location, allow_redirects=True, timeout=30)
        return jsonify(res.json())
    except Exception as e:
        return jsonify({'ok': False, 'error': str(e)})

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=80, debug=False)

